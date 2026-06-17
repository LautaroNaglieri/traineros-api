# -*- coding: utf-8 -*-
"""
Presentación ejecutiva TFI TrainerOS — las 5 fases completas, pensada para 10 minutos.
Estilo corporativo tipo "Naturgy en un Scan": azul #1F4E79 + naranja #E8741E + gris,
portada azul, títulos azules, tarjetas, acentos naranjas y barra naranja al pie.
Diagramas: imágenes draw.io ajustadas (sin distorsión).

Reparto (6 personas · 10:00):
  Conrado Gómez   — Intro (problema + negocio) y Cierre   ......... 2:30
  Nicolás Zyngale — Fase 1 (monolito en 3 capas) .................. 1:30
  Matías Saravia  — Fase 2 (microservicios A + B) ................. 1:30
  Lautaro Naglieri— Fase 3 (robustez y escalabilidad) ............. 1:30
  Pablo Czurylo   — Fase 4 (infra, DevOps y costos) ............... 1:30
  Leandro Gramajo — Fase 5 (integración con IA) ................... 1:30
"""
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BASE = r"C:\Users\Conrado\Desktop\Programacion\UNSTA\app-development\Final\Presentacion"
DIAG = os.path.join(BASE, "diagramas")

NAVY="1F4E79"; NAVY2="2A5E8C"; ORANGE="E8741E"; ORANGE2="C75F12"
CARD="ECF0F3"; CARD2="E1E8EE"; INK="404040"; GRAY="7A7A7A"; LGRAY="9AA6B2"; WHITE="FFFFFF"
def rgb(h): return RGBColor.from_string(h)
SW, SH = 13.333, 7.5

prs = Presentation(); prs.slide_width=Inches(SW); prs.slide_height=Inches(SH)
def blank(): return prs.slides.add_slide(prs.slide_layouts[6])

def rect(s,x,y,w,h,fill=None,line=None,line_w=1.0,radius=0.06,shape=MSO_SHAPE.RECTANGLE):
    sh=s.shapes.add_shape(shape,Inches(x),Inches(y),Inches(w),Inches(h))
    if fill is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb=rgb(fill)
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb=rgb(line); sh.line.width=Pt(line_w)
    sh.shadow.inherit=False
    if shape==MSO_SHAPE.ROUNDED_RECTANGLE:
        try: sh.adjustments[0]=radius
        except Exception: pass
    return sh

def P(runs,**kw): d={"runs":runs}; d.update(kw); return d
def txt(s,x,y,w,h,paras,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,wrap=True):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=wrap; tf.vertical_anchor=anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    for i,p in enumerate(paras):
        para=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        para.alignment=p.get("align",align)
        if "space_after" in p: para.space_after=Pt(p["space_after"])
        if "space_before" in p: para.space_before=Pt(p["space_before"])
        if "line" in p: para.line_spacing=p["line"]
        for (t,size,bold,color) in p["runs"]:
            r=para.add_run(); r.text=t; r.font.size=Pt(size); r.font.bold=bold
            r.font.color.rgb=rgb(color); r.font.name="Segoe UI"
    return tb

def bottom_bar(s): rect(s,0,SH-0.13,SW,0.13,fill=ORANGE)
def page(s,num,speaker=None,tiempo=None):
    parts=[p for p in [num, (f"Expone: {speaker}" if speaker else None), tiempo] if p]
    t="   ·   ".join(parts)
    txt(s,6.5,SH-0.46,6.28,0.28,[P([(t,9.5,False,LGRAY)],align=PP_ALIGN.RIGHT)],align=PP_ALIGN.RIGHT)

def head(s,title,subtitle=None):
    txt(s,0.7,0.55,12.0,0.7,[P([(title,30,True,NAVY)])])
    if subtitle:
        txt(s,0.72,1.28,12.0,0.4,[P([(subtitle,14.5,False,GRAY)])])

def fit_image(s,path,ax,ay,aw,ah):
    iw,ih=Image.open(path).size; ar=iw/ih; box=aw/ah
    if ar>box: w=aw; h=aw/ar
    else: h=ah; w=ah*ar
    s.shapes.add_picture(path,Inches(ax+(aw-w)/2),Inches(ay+(ah-h)/2),Inches(w),Inches(h))

def diagram_slide(name, num, speaker=None, tiempo=None, title=None):
    s=blank()
    if title:
        txt(s,0.7,0.22,12.0,0.5,[P([(title,20,True,NAVY)])])
        fit_image(s,os.path.join(DIAG,name),0.25,0.78,12.83,6.25)
    else:
        fit_image(s,os.path.join(DIAG,name),0.25,0.22,12.83,6.9)
    bottom_bar(s); page(s,num,speaker,tiempo)
    return s

INTEG="Conrado Gómez   ·   Nicolás Zyngale   ·   Matías Saravia   ·   Lautaro Naglieri   ·   Pablo Czurylo   ·   Leandro Gramajo"

# ---------------------------------------------------- 1. Portada  (Conrado · intro)
s=blank(); rect(s,0,0,SW,SH,fill=NAVY)
txt(s,0.8,2.35,11.73,1.0,[P([("TrainerOS",56,True,WHITE)],align=PP_ALIGN.CENTER)],align=PP_ALIGN.CENTER)
txt(s,0.8,3.52,11.73,0.55,[P([("“Del Código a la Arquitectura: El Criterio del Desarrollador”",22,True,ORANGE)],align=PP_ALIGN.CENTER)],align=PP_ALIGN.CENTER)
txt(s,0.8,4.18,11.73,0.45,[P([("Plataforma SaaS de gestión para entrenadores personales",15,False,"D6E0EA")],align=PP_ALIGN.CENTER)],align=PP_ALIGN.CENTER)
rect(s,0,4.85,SW,0.045,fill=ORANGE)
txt(s,0.8,5.4,11.73,0.4,[P([("Evolución arquitectónica en 5 fases: del monolito modular a los microservicios event-driven, la escala y la IA",12.5,False,"AEBFCE")],align=PP_ALIGN.CENTER)],align=PP_ALIGN.CENTER)
txt(s,0.8,6.35,11.73,0.6,[
    P([("TFI · Desarrollo de Aplicaciones Web · Grupo VOX",12.5,True,WHITE)],align=PP_ALIGN.CENTER,space_after=4),
    P([(INTEG,10.5,False,"AEBFCE")],align=PP_ALIGN.CENTER),
])

# ---------------------------------------------------- 2. Problema y negocio  (Conrado · 1:00)
s=blank()
head(s,"El problema y el modelo de negocio","Por qué existe TrainerOS y cómo gana dinero")
rect(s,0.7,1.95,5.95,2.4,fill=CARD)
txt(s,1.0,2.2,5.4,2.0,[
    P([("EL PROBLEMA",13,True,NAVY)],space_after=8),
    P([("Los entrenadores gestionan decenas de alumnos con ",13.5,False,INK),
       ("planillas de Excel y mensajes de WhatsApp",13.5,True,ORANGE),(".",13.5,False,INK)],space_after=7,line=1.18),
    P([("Caos logístico, rutinas desactualizadas y pérdida de control sobre los cobros de membresías.",13.5,False,INK)],line=1.18),
])
rect(s,6.85,1.95,5.78,2.4,fill=CARD)
txt(s,7.15,2.2,5.25,2.0,[
    P([("MODELO DE NEGOCIO — SaaS B2B2C",13,True,NAVY)],space_after=8),
    P([("Suscripción mensual por entrenador",13.5,True,ORANGE),
       (" (multitenant). El entrenador paga; sus alumnos usan gratis un panel móvil.",13.5,False,INK)],space_after=7,line=1.18),
    P([("Cada entrenador es un tenant aislado: sus datos nunca se cruzan con los de otro.",13.5,False,INK)],line=1.18),
])
txt(s,0.7,4.62,8,0.35,[P([("Funcionalidades del MVP",15,True,NAVY)])])
mvp=[("Auth con Google","Identidad y acceso"),("Constructor de rutinas","Sobre catálogo de ejercicios"),
     ("Panel móvil del alumno","Rutina del día y progreso"),("Control de pagos","Membresías y alertas")]
bx=0.7; bw=2.96; gap=0.097
for i,(t,dd) in enumerate(mvp):
    x=bx+i*(bw+gap)
    rect(s,x,5.1,bw,1.35,fill=CARD)
    rect(s,x,5.1,0.07,1.35,fill=ORANGE)
    txt(s,x+0.24,5.34,bw-0.42,0.95,[
        P([(t,12.5,True,NAVY)],space_after=4,line=1.05),
        P([(dd,10.5,False,GRAY)],line=1.1)])
bottom_bar(s); page(s,"2 / 11","Conrado Gómez","1:00")

# ---------------------------------------------------- 3. Fase 1 diagrama  (Nicolás · 0:50)
diagram_slide("fase1_monolito_3capas.png","3 / 11","Nicolás Zyngale","0:50")

# ---------------------------------------------------- 4. Fase 1 justificación  (Nicolás · 0:40)
s=blank()
head(s,"¿Por qué alcanza para arrancar? ¿Cuáles son sus límites?","El monolito como decisión deliberada, no como pecado")
rect(s,0.7,1.95,5.9,4.55,fill=CARD)
rect(s,0.7,1.95,5.9,0.09,fill=NAVY)
txt(s,1.0,2.2,5.35,0.4,[P([("POR QUÉ ALCANZA PARA ARRANCAR",13.5,True,NAVY)])])
just=[("Costo y time-to-market mínimos","Un deploy, una base y un pipeline. USD 7–16/mes; sprints semanales."),
      ("Un solo lenguaje de punta a punta","Node.js comparte ecosistema con el frontend (TypeScript full-stack)."),
      ("Ideal para carga I/O-bound","95% lecturas/escrituras, sin cómputo pesado: Express + Prisma rinde."),
      ("Modularidad intencional","Los módulos ya son los futuros microservicios. Abarata la Fase 2.")]
yy=2.78
for t,dd in just:
    txt(s,1.0,yy,5.3,0.9,[
        P([("■  ",10.5,True,NAVY),(t,12,True,INK)],space_after=2,line=1.05),
        P([("      "+dd,10.8,False,GRAY)],line=1.12)]); yy+=0.92
rect(s,6.75,1.95,5.88,4.55,fill=CARD)
rect(s,6.75,1.95,5.88,0.09,fill=ORANGE)
txt(s,7.05,2.2,5.35,0.4,[P([("SUS LÍMITES  (disparan la Fase 2)",13.5,True,ORANGE2)])])
lims=[("Solo escala vertical","Un pico en Rutinas obliga a escalar todo el monolito."),
      ("Acoplamiento de despliegue","Un deploy defectuoso de un módulo tumba todo el sistema."),
      ("Base única = riesgo concentrado","Una query pesada de reportes degrada lo transaccional."),
      ("Cuello de botella de equipos","Varios equipos: merge conflicts y releases que coordinar."),
      ("Sin alta disponibilidad","Render/Railway no dan autoscaling granular ni SLAs.")]
yy=2.78
for t,dd in lims:
    txt(s,7.05,yy,5.3,0.72,[
        P([("■  ",10.5,True,ORANGE),(t,12,True,INK)],space_after=1,line=1.05),
        P([("      "+dd,10.8,False,GRAY)],line=1.1)]); yy+=0.735
bottom_bar(s); page(s,"4 / 11","Nicolás Zyngale","0:40")

# ---------------------------------------------------- 5. Contexto Fase 2  (Matías · 0:30)
s=blank()
head(s,"El negocio fuerza la evolución a microservicios","12 meses después: convenios B2B con dos cadenas regionales de gimnasios")
rect(s,0.7,1.9,11.93,0.95,fill=NAVY)
txt(s,1.0,1.9,5.0,0.95,[
    P([("De ",13,False,"AEBFCE"),("100 entrenadores",18,True,WHITE)],space_after=1),
    P([("monolito · ~5.000 alumnos",11,False,"AEBFCE")])],anchor=MSO_ANCHOR.MIDDLE)
txt(s,6.0,1.9,1.0,0.95,[P([("→",26,True,ORANGE)],align=PP_ALIGN.CENTER)],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
txt(s,7.0,1.9,5.4,0.95,[
    P([("A ",13,False,"AEBFCE"),("5.000 entrenadores",18,True,ORANGE)],space_after=1),
    P([("~120.000 alumnos activos · equipos independientes",11,False,"AEBFCE")])],anchor=MSO_ANCHOR.MIDDLE)
txt(s,0.7,3.1,8,0.35,[P([("Tres presiones que el monolito no resuelve",15,True,NAVY)])])
press=[("Picos de tráfico violentos","De 18 a 21 h los alumnos consultan su rutina desde el gimnasio (lecturas masivas). Los lunes el pico se triplica."),
       ("Equipos independientes","Tres equipos (Core de rutinas, Pagos/B2B, Experiencia del alumno) que necesitan desplegar sin coordinarse."),
       ("Requisitos B2B","Las cadenas exigen SLA, facturación consolidada y notificaciones masivas que no pueden competir con el flujo crítico.")]
px0=0.7; pw=3.94; pg=0.155
for i,(t,dd) in enumerate(press):
    x=px0+i*(pw+pg)
    rect(s,x,3.6,pw,2.7,fill=CARD)
    rect(s,x+0.28,3.85,0.6,0.6,fill=ORANGE)
    txt(s,x+0.28,3.85,0.6,0.6,[P([(str(i+1),20,True,WHITE)],align=PP_ALIGN.CENTER)],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,x+0.28,4.62,pw-0.56,1.55,[
        P([(t,14,True,NAVY)],space_after=5,line=1.05),
        P([(dd,11,False,INK)],line=1.2)])
bottom_bar(s); page(s,"5 / 11","Matías Saravia","0:30")

# ---------------------------------------------------- 6,7 diagramas Fase 2  (Matías · 0:30 c/u)
diagram_slide("fase2a_microservicios_tradicionales.png","6 / 11","Matías Saravia","0:30")
diagram_slide("fase2b_microservicios_modernos.png","7 / 11","Matías Saravia","0:30")

# ---------------------------------------------------- 8. Fase 3 diagrama  (Lautaro · 1:30)
diagram_slide("fase3_arquitectura_escalada.png","8 / 11","Lautaro Naglieri","1:30")

# ---------------------------------------------------- 9. Fase 4 infra/DevOps/costos  (Pablo · 1:30)
s=blank()
head(s,"Infraestructura, DevOps y ambientes","Google Cloud Platform — un mismo artefacto Docker promovido por ambiente")
quad=[("AMBIENTES (4 proyectos GCP)",
       ["4 proyectos aislados: Desarrollo · QA · UAT · Producción (permisos, red y facturación separados).",
        "Se promueve la MISMA imagen Docker (tag = commit SHA); solo cambia la config inyectada."]),
      ("CI/CD — GitHub Actions",
       ["CI en cada PR: lint → tests → build → scan Trivy → push a Artifact Registry.",
        "CD: QA automático · UAT (rc-*) · Prod con aprobación manual y rollback por healthchecks."]),
      ("CONTENEDORES + KUBERNETES (GKE)",
       ["1 imagen distroless (~80 MB) por servicio; Deployments con probes y HPA.",
        "Ingress → Cloud Load Balancer; secretos con Secret Manager + Workload Identity."]),
      ("INFRAESTRUCTURA COMO CÓDIGO (Terraform)",
       ["Módulos reutilizables: network · gke · cloudsql · redis · pubsub · iam, con dir por ambiente.",
        "Estado remoto en GCS; cada cambio sigue el flujo PR → plan → apply."])]
qx=[0.7,6.85]; qy=[1.95,4.05]; qw=5.78; qh=1.95
accent=[NAVY,ORANGE,NAVY,ORANGE]
for i,(t,items) in enumerate(quad):
    x=qx[i%2]; y=qy[i//2]
    rect(s,x,y,qw,qh,fill=CARD)
    rect(s,x,y,0.08,qh,fill=accent[i])
    paras=[P([(t,12.5,True,NAVY)],space_after=6,line=1.05)]
    for it in items:
        paras.append(P([("•  ",11,True,accent[i]),(it,11,False,INK)],space_after=4,line=1.12))
    txt(s,x+0.28,y+0.2,qw-0.5,qh-0.3,paras)
rect(s,0.7,6.22,11.93,0.62,fill=NAVY)
txt(s,1.0,6.22,11.4,0.62,[P([
    ("Costos (producción, GCP): ",11.5,True,ORANGE),
    ("~USD 810/mes total · QA+UAT ~USD 150/mes (nodos spot, apagado nocturno) ≈ 1,1% del MRR con 5.000 entrenadores × USD 15.",11.5,False,WHITE)])],anchor=MSO_ANCHOR.MIDDLE)
bottom_bar(s); page(s,"9 / 11","Pablo Czurylo","1:30")

# ---------------------------------------------------- 10. Fase 5 diagrama  (Leandro · 1:30)
diagram_slide("fase5_integracion_ia.png","10 / 11","Leandro Gramajo","1:30")

# ---------------------------------------------------- 11. Cierre / evolución  (Conrado · 1:00)
s=blank()
head(s,"La evolución en una mirada","Cada fase resolvió un problema concreto del negocio — nunca tecnología “porque sí”")
fases=[("FASE 1","Monolito modular en 3 capas","MVP barato y rápido: un deploy, una base. USD 7–16/mes.",NAVY),
       ("FASE 2","Microservicios: REST → eventos","Autonomía de equipos; Pub/Sub desacopla y tolera fallos.",ORANGE),
       ("FASE 3","Robustez y escalabilidad","Redis, read replicas, LB y auto-scaling absorben el pico.",NAVY),
       ("FASE 4","Infraestructura y DevOps","GKE + Terraform + CI/CD; 4 ambientes. ~1,1% del MRR.",ORANGE),
       ("FASE 5","Integración con IA","LLM asíncrono sobre el backbone de eventos; RAG a futuro.",NAVY)]
yy=1.95; rh=0.78
for tag,t,dd,col in fases:
    rect(s,0.7,yy,1.7,rh-0.08,fill=col)
    txt(s,0.7,yy,1.7,rh-0.08,[P([(tag,13,True,WHITE)],align=PP_ALIGN.CENTER)],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    rect(s,2.5,yy,10.13,rh-0.08,fill=CARD)
    txt(s,2.78,yy,9.7,rh-0.08,[
        P([(t,13.5,True,NAVY)],space_after=2,line=1.0),
        P([(dd,11.5,False,INK)],line=1.05)],anchor=MSO_ANCHOR.MIDDLE)
    yy+=rh
rect(s,0.7,yy+0.04,11.93,0.66,fill=ORANGE)
txt(s,1.0,yy+0.04,11.4,0.66,[P([
    ("Conclusión: ",12,True,WHITE),
    ("cada fase respondió a una necesidad real del negocio y la arquitectura escaló de forma incremental, justificada por costo y rendimiento: el criterio del desarrollador por encima de la tecnología.",12,False,WHITE)])],anchor=MSO_ANCHOR.MIDDLE)
bottom_bar(s); page(s,"11 / 11","Conrado Gómez","1:00")

out=os.path.join(BASE,"TrainerOS_TFI_Presentacion.pptx")
prs.save(out)
print("OK ->",out,"| slides:",len(prs.slides._sldIdLst))
